# report_generator.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import io
import plotly.io as pio
from datetime import date

def add_title_slide(pres, title_text, subtitle_text):
    """Adds a title slide to the presentation."""
    slide_layout = pres.slide_layouts[0]
    slide = pres.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def add_kpi_slide(pres, kpi_title, kpi_data):
    """Adds a slide with KPI metrics."""
    slide_layout = pres.slide_layouts[5] # Title Only layout
    slide = pres.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = kpi_title
    
    # Define positions for 4 KPIs
    positions = [
        {'left': Inches(0.5), 'top': Inches(1.5), 'width': Inches(2.0), 'height': Inches(1.5)},
        {'left': Inches(2.8), 'top': Inches(1.5), 'width': Inches(2.0), 'height': Inches(1.5)},
        {'left': Inches(5.1), 'top': Inches(1.5), 'width': Inches(2.0), 'height': Inches(1.5)},
        {'left': Inches(7.4), 'top': Inches(1.5), 'width': Inches(2.0), 'height': Inches(1.5)},
    ]

    for i, (metric_name, metric_value) in enumerate(kpi_data.items()):
        if i < len(positions):
            pos = positions[i]
            textbox = slide.shapes.add_textbox(pos['left'], pos['top'], pos['width'], pos['height'])
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            p_name = text_frame.add_paragraph()
            p_name.text = metric_name
            p_name.font.size = Pt(16)
            p_name.font.bold = True
            
            p_value = text_frame.add_paragraph()
            p_value.text = str(metric_value)
            p_value.font.size = Pt(28)
            p_value.font.color.rgb = RGBColor(0, 57, 166) # QuidelOrtho Blue

def add_chart_slide(pres, chart_title, fig, notes=""):
    """Adds a slide with a Plotly chart and optional notes."""
    slide_layout = pres.slide_layouts[6] # Blank layout
    slide = pres.slides.add_slide(slide_layout)

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    p = title_frame.add_paragraph()
    p.text = chart_title
    p.font.size = Pt(24)
    p.font.bold = True

    # Convert Plotly fig to image and add to slide
    img_bytes = pio.to_image(fig, format="png", width=800, height=450, scale=2)
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, Inches(1), Inches(1.2), width=Inches(8))

    # Add notes
    if notes:
        notes_shape = slide.shapes.add_textbox(Inches(1), Inches(6.0), Inches(8), Inches(1.2))
        notes_frame = notes_shape.text_frame
        p_notes = notes_frame.add_paragraph()
        p_notes.text = "Manager's Analysis:"
        p_notes.font.bold = True
        p_notes.font.size = Pt(14)

        p_notes_body = notes_frame.add_paragraph()
        p_notes_body.text = notes
        p_notes_body.font.size = Pt(12)

def generate_monthly_review_ppt(kpi_data, fig_timeline, fig_risk):
    """Generates the full PowerPoint presentation."""
    pres = Presentation()
    pres.slide_width = Inches(10)
    pres.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(pres, "Validation Monthly Management Review", f"Status as of: {date.today().strftime('%B %d, %Y')}")

    # Slide 2: KPIs
    add_kpi_slide(pres, "Executive Validation Portfolio Health", kpi_data)

    # Slide 3: Timeline Chart
    add_chart_slide(pres, "Portfolio Timeline & Critical Path", fig_timeline,
                    notes="Critical path projects (red outline) represent the highest risk to overall program timelines. "
                          "Any delays in these projects require immediate escalation and mitigation planning.")
    
    # Slide 4: Risk Chart
    add_chart_slide(pres, "Integrated Risk Posture (ISO 14971)", fig_risk,
                    notes="Focus remains on mitigating the high-score risks in the 'Unacceptable Region'. "
                          "The histogram distribution shows a concentration of 'Probable' risks, suggesting a need for more robust upstream design controls.")

    # Save to a byte stream
    ppt_stream = io.BytesIO()
    pres.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

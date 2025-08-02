# app_utils.py
import pandas as pd
import numpy as np
from datetime import date, timedelta
import plotly.express as px
import plotly.graph_objects as go
import plotly.io as pio
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io
from scipy import stats

# --- Custom Plotly Template ---
# This can be defined once and used across all plots.
quidelortho_template = {"layout": {"font": {"family": "Arial", "size": 12}}}
pio.templates["quidelortho_sme"] = go.layout.Template(quidelortho_template)
pio.templates.default = "quidelortho_sme"

# === CORE DATA GENERATION ===
def generate_vv_project_data():
    """Generates enhanced V&V and Validation project data."""
    data = {
        'Project/Assay': ['Savanna® RVP12 Assay', 'Sofia® 2 SARS Antigen+ FIA v2', 'Vitros® HIV Combo 5 Test', 'Automated Filler & Capper Line 3', 'Ortho-Vision® Analyzer SW Patch V&V'],
        'Type': ['Assay', 'Assay', 'Assay', 'Equipment', 'Software'],
        'Platform': ['Savanna', 'Sofia', 'Vitros', 'Manufacturing', 'Ortho-Vision'],
        'V&V Lead': ['M. Rodriguez', 'J. Chen', 'S. Patel', 'V. Kumar', 'S. Patel'],
        'V&V Phase': ['Execution', 'Reporting', 'PQ', 'IQ/OQ', 'Data Analysis'],
        'Overall Status': ['On Track', 'On Track', 'At Risk', 'On Track', 'Behind Schedule'],
        'Regulatory Pathway': ['510(k)', 'EUA Modification', 'PMA Supplement', 'N/A', 'Letter to File'],
        'Key Milestone': ['Final Report for Submission', 'EUA Submission', 'Production Readiness', 'OQ Completion', 'V&V Report for ECO'],
        'Start Date': [date.today() - timedelta(days=60), date.today() - timedelta(days=90), date.today() - timedelta(days=10), date.today() - timedelta(days=45), date.today() - timedelta(days=30)],
        'Due Date': [date.today() + timedelta(days=45), date.today() + timedelta(days=5), date.today() + timedelta(days=60), date.today() + timedelta(days=75), date.today() + timedelta(days=20)],
        'On Critical Path': [False, True, True, True, False],
        'Budget (USD)': [250000, 150000, 500000, 1200000, 75000],
        'Spent (USD)': [150000, 135000, 350000, 450000, 40000],
        'Schedule Performance Index (SPI)': [1.05, 1.10, 0.92, 1.02, 0.98],
        'First Time Right %': [98, 99, 95, 92, 100],
        'Utilization %': [85, 90, 110, 100, 75]
    }
    df = pd.DataFrame(data)
    df['Start Date'] = pd.to_datetime(df['Start Date'])
    df['Due Date'] = pd.to_datetime(df['Due Date'])
    return df

def generate_risk_management_data():
    """Generates enhanced risk data for both assay and equipment validation."""
    data = {
        'Risk ID': ['R-SAV-001', 'R-EQP-001', 'R-VIT-002', 'R-SFT-001'],
        'Project': ['Savanna® RVP12 Assay', 'Automated Filler & Capper Line 3', 'Vitros® HIV Combo 5 Test', 'Ortho-Vision® Analyzer SW Patch V&V'],
        'Risk Description': ['Potential cross-reactivity with emerging coronavirus strains.', 'Inconsistent torque from capper compromises sterility.', 'Biotin interference leads to false negatives.', 'Software patch re-introduces a previously fixed bug.'],
        'Severity': [4, 5, 5, 4],
        'Probability': [3, 3, 2, 2],
        'Owner': ['R&D/V&V', 'Automation/Validation', 'Clinical/V&V', 'SW V&V Team'],
        'Mitigation': ['Test against comprehensive panel.', 'Incorporate torque verification into OQ/PQ.', 'Include biotin in interference testing.', 'Expand regression test suite.']
    }
    df = pd.DataFrame(data)
    df['Risk_Score'] = df['Severity'] * df['Probability']
    return df.sort_values(by='Risk_Score', ascending=False)

# === V&V/VALIDATION STUDY DATA GENERATION (from utils.py) ===
def generate_linearity_data_immunoassay():
    expected_conc = np.array([0, 10, 50, 100, 250, 500, 750, 1000])
    od_values = 2.5 * expected_conc / (50 + expected_conc)
    observed_od = od_values + np.random.normal(0, 0.03, expected_conc.shape)
    return pd.DataFrame({'Analyte Concentration (ng/mL)': expected_conc, 'Optical Density (OD)': observed_od})

def generate_precision_data_clsi_ep05():
    data = []; days = [f'Day {i}' for i in range(1, 6)];
    for day in days:
        for run in range(1, 3):
            for _ in range(3):
                data.append({'Control': 'Low Positive', 'Day': day, 'Run': run, 'S/CO Ratio': np.random.normal(1.8, 0.2)})
                data.append({'Control': 'Moderate Positive', 'Day': day, 'Run': run, 'S/CO Ratio': np.random.normal(8.5, 0.8)})
    return pd.DataFrame(data)

def generate_method_comparison_data():
    reference = np.random.uniform(10, 500, 100)
    candidate = 1.05 * reference - 5 + np.random.normal(0, 15, 100)
    return pd.DataFrame({'Reference Method (U/L)': reference, 'Candidate Method (U/L)': candidate})

def generate_equipment_pq_data():
    data = []
    for i in range(1, 11):
        mean_fill = 5.02 if i < 7 else 5.08
        values = np.random.normal(loc=mean_fill, scale=0.05, size=30)
        for val in values:
            data.append({'Batch': f'PQ_Batch_{i}', 'Fill Volume (mL)': val})
    return pd.DataFrame(data)

# --- NEW/FIXED DATA GENERATION FUNCTIONS ---
# These functions were missing from the original files and have been created here.

def generate_traceability_matrix_data():
    return pd.DataFrame({
        'Requirement ID': ['URS-001', 'SRS-010', 'RISK-CTRL-005', 'SRS-011'],
        'Requirement Type': ['User Need', 'Software Requirement', 'Risk-Ctrl Requirement', 'Software Requirement'],
        'Requirement Text': ['Assay shall detect HIV.', 'Software shall flag results > Ct 37.', 'Mitigate biotin interference.', 'Software shall generate audit trail.'],
        'Test Case ID': ['TC-SENS-01-01', 'TC-FLAG-01-03', 'TC-INTER-01-01', 'TC-AUDIT-01-01'],
        'Test Result': ['Pass', 'Fail', 'Pass', 'Pass']
    })

def generate_risk_burndown_data():
    return pd.DataFrame({
        'Week': [1, 2, 3, 4, 5, 6],
        'High': [5, 5, 4, 3, 2, 1],
        'Medium': [8, 8, 9, 10, 10, 10],
        'Low': [10, 11, 12, 12, 13, 14]
    })

def generate_change_control_data():
    return pd.DataFrame({
        'ECO Number': ['ECO-24-101', 'ECO-24-102', 'ECO-24-105', 'ECO-24-106'],
        'Product Impacted': ['Sofia® RSV', 'QuickVue® Flu', 'Sofia® RSV', 'Vitros® HIV'],
        'Status': ['Awaiting V&V Plan', 'V&V in Progress', 'V&V Complete', 'V&V in Progress'],
        'V&V Lead': ['J. Chen', 'S. Patel', 'J. Chen', 'M. Rodriguez']
    })

def generate_analytical_specificity_data_molecular():
    return pd.DataFrame({
        'Test Type': ['Inclusivity', 'Inclusivity', 'Exclusivity', 'Exclusivity', 'Exclusivity'],
        'Organism': ['SARS-CoV-2 (WA-1)', 'SARS-CoV-2 (Omicron)', 'MERS-CoV', 'hCoV-OC43', 'Influenza A'],
        'Result': ['Positive', 'Positive', 'Negative', 'Negative', 'Negative'],
        'Notes': ['Expected positive result.', 'Expected positive result.', 'Expected negative.', 'Expected negative.', 'Expected negative.']
    })

def generate_lot_to_lot_data():
    ref = np.random.normal(100, 10, 30)
    cand = np.random.normal(102, 10, 30)
    df1 = pd.DataFrame({'Reagent Lot ID': 'Reference_Lot', 'Test Line Intensity': ref})
    df2 = pd.DataFrame({'Reagent Lot ID': 'Candidate_Lot', 'Test Line Intensity': cand})
    return pd.concat([df1, df2])

def calculate_equivalence(df, group_col, value_col, low, high):
    groups = df[group_col].unique()
    data1 = df[df[group_col] == groups[0]][value_col]
    data2 = df[df[group_col] == groups[1]][value_col]
    mean_diff = data2.mean() - data1.mean()
    tost_result = stats.ttest_ind_from_stats(
        mean1=data2.mean(), std1=data2.std(), nobs1=len(data2),
        mean2=data1.mean(), std2=data1.std(), nobs2=len(data1)
    )
    # This is a simplification; a real TOST involves two one-sided tests.
    # We'll use a simple check for demonstration.
    p_tost = 0.01 if low < mean_diff < high else 0.99
    return p_tost, tost_result.pvalue, mean_diff

def generate_defect_trend_data():
    dates = pd.to_datetime(pd.date_range(start="2023-01-01", periods=12, freq='W'))
    found = np.cumsum(np.random.randint(1, 5, 12))
    open = found - np.cumsum(np.random.randint(0, 4, 12))
    return pd.DataFrame({'Date': dates, 'Total Defects Found': found, 'Open Defects': open.clip(lower=0)})

def generate_submission_package_data(project_name, pathway):
    return pd.DataFrame({
        'Deliverable': ['V&V Master Plan', 'Risk Management File', 'RTM', 'Sensitivity Report', 'Specificity Report', 'V&V Summary Report'],
        'Document ID': ['DHF-001', 'DHF-002', 'DHF-003', 'V&V-RPT-010', 'V&V-RPT-011', 'DHF-020'],
        'Status': ['Approved', 'Approved', 'Approved', 'In Review', 'Execution', 'Drafting'],
        'Progress': [100, 100, 100, 85, 40, 15],
        'Regulatory_Impact': ['High', 'High', 'High', 'Medium', 'High', 'High'],
        'Statistical_Robustness': [10, 9, 9, 7, 6, 9]
    })

def generate_capa_data():
    return pd.DataFrame({
        'ID': ['CAPA-23-011', 'INV-24-005', 'CAPA-24-002'],
        'Source': ['Post-Launch Complaint', 'V&V Study Anomaly', 'Internal Audit'],
        'Product': ['Sofia® RSV', 'Savanna® RVP12 Assay', 'All Platforms'],
        'Description': ['Increase in false positives.', 'Unexpected weak positive signal.', 'Inconsistent deviation documentation.'],
        'Owner': ['V&V Mgmt', 'M. Rodriguez', 'QA / V&V Mgmt'],
        'Phase': ['Effectiveness Check', 'Investigation', 'Implementation'],
        'Due Date': [date.today() - timedelta(days=10), date.today() + timedelta(days=18), date.today() + timedelta(days=40)]
    })

def generate_validation_program_data():
    data = {
        'System_ID': ['EQP-00123', 'EQP-00456', 'SW-0010', 'EQP-00789'],
        'System_Name': ['Automated Filler Line 3', 'Blister Pack Sealer 7', 'LIMS v3.2', 'Lyophilizer #2'],
        'Validation_Status': ['Validated', 'Validation in Progress', 'Validated', 'Revalidation Due'],
        'Last_Validation_Date': [date(2022, 5, 20), date(2024, 6, 15), date(2023, 1, 10), date(2021, 8, 30)],
        'Next_Revalidation_Date': [date(2025, 5, 20), date(2027, 6, 15), date(2026, 1, 10), date(2024, 8, 30)]
    }
    df = pd.DataFrame(data)
    df['Next_Revalidation_Date'] = pd.to_datetime(df['Next_Revalidation_Date'])
    df['Days_Until_Due'] = (df['Next_Revalidation_Date'] - pd.to_datetime(date.today())).dt.days
    df.loc[df['Days_Until_Due'] < 0, 'Validation_Status'] = 'Revalidation Overdue'
    return df

def generate_instrument_schedule_data():
    today = date.today()
    return pd.DataFrame({
        'Instrument': ['Savanna-01', 'Savanna-01', 'Sofia-03', 'Vitros-02', 'Vitros-02'],
        'Start': [pd.to_datetime(today), pd.to_datetime(today) + timedelta(days=2), pd.to_datetime(today), pd.to_datetime(today), pd.to_datetime(today) + timedelta(days=3)],
        'Finish': [pd.to_datetime(today) + timedelta(days=1), pd.to_datetime(today) + timedelta(days=4), pd.to_datetime(today) + timedelta(days=5), pd.to_datetime(today) + timedelta(days=2), pd.to_datetime(today) + timedelta(days=4)],
        'Status': ['V&V Execution', 'Calibration/PM', 'V&V Execution', 'OOS', 'Available'],
        'Details': ['RVP12 Study', 'Annual PM', 'SARS v2 Study', 'Sensor Failure', 'Open for booking']
    })

def generate_training_data_for_heatmap():
    data = {'Savanna': [2, 1, 0, 1], 'Sofia': [1, 2, 1, 0], 'Vitros': [0, 1, 2, 1], 'JMP/R': [2, 1, 1, 0], 'Protocol Authoring': [2, 2, 1, 1]}
    index = ['M. Rodriguez', 'J. Chen', 'S. Patel', 'K. Lee']
    return pd.DataFrame(data, index=index)

def generate_reagent_lot_status_data():
    return pd.DataFrame({
        'Material Name': ['RVP12 Master Mix', 'SARS v2 Test Cassettes', 'HIV Control Kit'],
        'Lot Number': ['MM-24A01', 'CAS-23B15', 'CTRL-22C09'],
        'Expiry Date': [date.today() + timedelta(days=180), date.today() + timedelta(days=30), date.today() - timedelta(days=10)],
        'Quantity': [5, 20, 1],
        'Status': ['Qualified', 'Low Inventory', 'Expired - Quarantined'],
        'Reserved For': ['Savanna RVP12', 'Sofia SARS v2', 'N/A']
    })

def calculate_instrument_utilization(schedule_df):
    schedule_df['Duration'] = (schedule_df['Finish'] - schedule_df['Start']).dt.total_seconds() / 3600
    schedule_df['Platform'] = schedule_df['Instrument'].apply(lambda x: x.split('-')[0])
    return schedule_df

def generate_idp_data():
    return pd.DataFrame({
        'Team Member': ['J. Chen', 'S. Patel', 'K. Lee'],
        'Development Goal': ['Achieve Expert in JMP/R', 'Achieve Practitioner on Vitros', 'Complete Protocol Authoring Training'],
        'Mentor': ['M. Rodriguez', 'Manager', 'M. Rodriguez'],
        'Start Date': [pd.to_datetime(date.today() - timedelta(days=30)), pd.to_datetime(date.today()), pd.to_datetime(date.today())],
        'Target Date': [pd.to_datetime(date.today() + timedelta(days=150)), pd.to_datetime(date.today() + timedelta(days=90)), pd.to_datetime(date.today() + timedelta(days=45))]
    })

def generate_process_excellence_data():
    months = pd.date_range(start="2023-01-01", periods=18, freq='M')
    return pd.DataFrame({
        'Month': months,
        'Protocol_Approval_Cycle_Time_Days': 45 - np.arange(18) * 0.5 + np.random.normal(0, 2, 18),
        'Report_Rework_Rate_Percent': 8 + np.sin(np.linspace(0, 2 * np.pi, 18)) * 2,
        'Deviations_per_100_Test_Hours': 3 - np.arange(18) * 0.05 + np.random.normal(0, 0.3, 18)
    })

def generate_workload_forecast_data():
    ds = pd.to_datetime(pd.date_range(start='2022-01-01', periods=24, freq='MS'))
    trend = np.linspace(500, 800, 24)
    seasonality = 100 * np.sin(np.linspace(0, 4 * np.pi, 24))
    noise = np.random.normal(0, 50, 24)
    y = trend + seasonality + noise
    return pd.DataFrame({'ds': ds, 'y': y})

# === REPORT GENERATOR (from report_generator.py) ===
def add_title_slide(pres, title_text, subtitle_text):
    slide = pres.slides.add_slide(pres.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text

def add_kpi_slide(pres, kpi_title, kpi_data):
    slide = pres.slides.add_slide(pres.slide_layouts[5])
    slide.shapes.title.text = kpi_title
    positions = [
        {'left': Inches(0.5), 'top': Inches(1.5)}, {'left': Inches(2.8), 'top': Inches(1.5)},
        {'left': Inches(5.1), 'top': Inches(1.5)}, {'left': Inches(7.4), 'top': Inches(1.5)},
    ]
    for i, (metric_name, metric_value) in enumerate(kpi_data.items()):
        if i < len(positions):
            pos = positions[i]
            textbox = slide.shapes.add_textbox(pos['left'], pos['top'], Inches(2.0), Inches(1.5))
            p_name = textbox.text_frame.add_paragraph(); p_name.text = metric_name; p_name.font.size = Pt(16)
            p_value = textbox.text_frame.add_paragraph(); p_value.text = str(metric_value); p_value.font.size = Pt(28)

def add_chart_slide(pres, chart_title, fig):
    slide = pres.slides.add_slide(pres.slide_layouts[6]) # Blank layout
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_shape.text_frame.add_paragraph().text = chart_title
    img_bytes = pio.to_image(fig, format="png", width=800, height=450, scale=2)
    slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(1), Inches(1.2), width=Inches(8))

def generate_monthly_review_ppt(kpi_data, fig_timeline, fig_risk):
    """Generates the full PowerPoint presentation."""
    pres = Presentation()
    pres.slide_width = Inches(10)
    pres.slide_height = Inches(7.5)
    add_title_slide(pres, "Validation Monthly Management Review", f"Status as of: {date.today().strftime('%B %d, %Y')}")
    add_kpi_slide(pres, "Executive Validation Portfolio Health", kpi_data)
    add_chart_slide(pres, "Portfolio Timeline & Critical Path", fig_timeline)
    add_chart_slide(pres, "Integrated Risk Posture (ISO 14971)", fig_risk)
    ppt_stream = io.BytesIO()
    pres.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

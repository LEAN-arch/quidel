# app_utils.py
import pandas as pd
import numpy as np
from datetime import date, timedelta
import io
from scipy import stats
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import plotly.io as pio

# --- Custom Plotly Template ---
pio.templates.default = "plotly_white"

# === CORE DATA GENERATION (Fully Restored from Original Files) ===

def generate_vv_project_data():
    """Generates enhanced V&V and Validation project data with financial and performance metrics."""
    data = {
        'Project/Assay': ['Savanna® RVP12 Assay', 'Sofia® 2 SARS Antigen+ FIA v2', 'Vitros® HIV Combo 5 Test', 'Triage® BNPNext™ Consumable Change', 'Automated Filler & Capper Line 3', 'Ortho-Vision® Analyzer SW Patch V&V', 'New Blister Pack Sealer Unit 7'],
        'Type': ['Assay', 'Assay', 'Assay', 'Assay', 'Equipment', 'Software', 'Equipment'],
        'Platform': ['Savanna', 'Sofia', 'Vitros', 'Triage', 'Manufacturing', 'Ortho-Vision', 'Packaging'],
        'V&V Lead': ['M. Rodriguez', 'J. Chen', 'S. Patel', 'M. Rodriguez', 'V. Kumar', 'S. Patel', 'V. Kumar'],
        'V&V Phase': ['Execution', 'Reporting', 'PQ', 'On Hold', 'IQ/OQ', 'Data Analysis', 'FAT/SAT'],
        'Overall Status': ['On Track', 'On Track', 'At Risk', 'On Hold', 'On Track', 'On Track', 'Behind Schedule'],
        'Regulatory Pathway': ['510(k)', 'EUA Modification', 'PMA Supplement', 'Letter to File', 'N/A', 'Letter to File', 'N/A'],
        'Key Milestone': ['Final Report for Submission', 'EUA Submission', 'Production Readiness', 'N/A', 'OQ Completion', 'V&V Report for ECO', 'SAT Completion'],
        'Start Date': [date.today() - timedelta(days=60), date.today() - timedelta(days=90), date.today() - timedelta(days=10), date.today() - timedelta(days=120), date.today() - timedelta(days=45), date.today() - timedelta(days=30), date.today() - timedelta(days=90)],
        'Due Date': [date.today() + timedelta(days=45), date.today() + timedelta(days=5), date.today() + timedelta(days=60), date.today() + timedelta(days=15), date.today() + timedelta(days=75), date.today() + timedelta(days=20), date.today() + timedelta(days=30)],
        'On Critical Path': [False, False, True, False, True, False, True],
        'Budget (USD)': [250000, 150000, 500000, 50000, 1200000, 75000, 300000],
        'Spent (USD)': [150000, 135000, 350000, 20000, 450000, 40000, 250000],
        'Schedule Performance Index (SPI)': [1.05, 1.10, 0.92, 1.0, 1.02, 0.98, 0.85],
        'First Time Right %': [98, 99, 95, 100, 92, 100, 88],
        'Utilization %': [85, 90, 110, 0, 100, 75, 120]
    }
    df = pd.DataFrame(data)
    df['Start Date'] = pd.to_datetime(df['Start Date'])
    df['Due Date'] = pd.to_datetime(df['Due Date'])
    return df

def generate_risk_management_data():
    """Generates enhanced risk data for both assay and equipment validation."""
    data = {
        'Risk ID': ['R-SAV-001', 'R-EQP-001', 'R-VIT-002', 'R-TRI-001', 'R-GEN-005', 'R-SFT-001'],
        'Project': ['Savanna® RVP12 Assay', 'Automated Filler & Capper Line 3', 'Vitros® HIV Combo 5 Test', 'Triage® BNPNext™ Consumable Change', 'Savanna® RVP12 Assay', 'Ortho-Vision® Analyzer SW Patch V&V'],
        'Risk Description': ['Potential cross-reactivity with emerging non-pathogenic coronavirus strains could lead to false positives.','New automated capper applies inconsistent torque, potentially compromising container closure integrity and product sterility.','Undetected interference from biotin could lead to a false negative result, delaying critical diagnosis.','New plastic supplier for consumable has higher lot-to-lot variability, potentially causing inconsistent test results.','V&V study execution timeline conflicts with key personnel availability, posing a risk of delayed regulatory submission.','Regression testing for software patch fails to cover a legacy edge-case, potentially re-introducing a previously fixed bug.'],
        'Severity': [4, 5, 5, 3, 2, 4], 'Probability': [3, 3, 2, 3, 4, 2],
        'Owner': ['R&D/V&V', 'Automation/Validation', 'Clinical/V&V', 'Supply Chain/V&V', 'V&V Management', 'SW V&V Team'],
        'Mitigation': ['Test against a comprehensive panel of endemic coronaviruses per FDA guidance.','Incorporate torque verification and container closure integrity testing (CCIT) into OQ and PQ protocols.','Include biotin and other interferents in testing per CLSI EP07. Add limitation to Instructions for Use (IFU).','Require tighter Certificate of Analysis (CoA) specs from supplier; perform incoming QC on each lot.','Re-allocate V&V specialist from a lower priority project. Document change in project plan.','Expand regression test suite to include historical defect scenarios.']
    }
    df = pd.DataFrame(data)
    df['Risk_Score'] = df['Severity'] * df['Probability']
    return df.sort_values(by='Risk_Score', ascending=False)

def generate_linearity_data_immunoassay():
    expected_conc = np.array([0, 10, 50, 100, 250, 500, 750, 1000])
    od_values = 2.5 * expected_conc / (50 + expected_conc); observed_od = od_values + np.random.normal(0, 0.03, expected_conc.shape)
    return pd.DataFrame({'Analyte Concentration (ng/mL)': expected_conc, 'Optical Density (OD)': observed_od})

def generate_precision_data_clsi_ep05():
    data = []; days = [f'Day {i}' for i in range(1, 6)];
    for day in days:
        for run in range(1, 3):
            for _ in range(3):
                data.append({'Control': 'Low Positive (Near Cutoff)', 'Day': day, 'Run': run, 'S/CO Ratio': np.random.normal(1.8, 0.2)})
                data.append({'Control': 'Moderate Positive', 'Day': day, 'Run': run, 'S/CO Ratio': np.random.normal(8.5, 0.8)})
    return pd.DataFrame(data)

def generate_method_comparison_data():
    ref = np.random.uniform(10, 500, 100); cand = 1.05 * ref - 5 + np.random.normal(0, 15, 100)
    return pd.DataFrame({'Reference Method (U/L)': ref, 'Candidate Method (U/L)': cand})

def generate_equipment_pq_data():
    data = []
    for i in range(1, 11):
        mean_fill = 5.02 if i < 7 else 5.08
        for val in np.random.normal(loc=mean_fill, scale=0.05, size=30): data.append({'Batch': f'PQ_Batch_{i}', 'Fill Volume (mL)': val})
    return pd.DataFrame(data)

def generate_traceability_matrix_data():
    return pd.DataFrame({'Requirement ID': ['URS-001', 'SRS-010', 'RISK-CTRL-005', 'SRS-011'], 'Requirement Type': ['User Need', 'Software Requirement', 'Risk-Ctrl Requirement', 'Software Requirement'], 'Test Case ID': ['TC-SENS-01', 'TC-FLAG-01', 'TC-INTER-01', 'TC-AUDIT-01'], 'Test Result': ['Pass', 'Fail', 'Pass', 'Pass']})

def generate_risk_burndown_data():
    return pd.DataFrame({'Week': [1, 2, 3, 4, 5, 6], 'High': [5, 5, 4, 3, 2, 1], 'Medium': [8, 8, 9, 10, 10, 10], 'Low': [10, 11, 12, 12, 13, 14]})

def generate_change_control_data():
    return pd.DataFrame({'ECO Number': ['ECO-24-101', 'ECO-24-102', 'ECO-24-105', 'ECO-24-106'], 'Product Impacted': ['Sofia® RSV', 'QuickVue® Flu', 'Sofia® RSV', 'Vitros® HIV'], 'Status': ['Awaiting V&V Plan', 'V&V in Progress', 'V&V Complete', 'V&V in Progress']})

def generate_analytical_specificity_data_molecular():
    return pd.DataFrame({'Test Type': ['Inclusivity', 'Exclusivity'], 'Organism': ['SARS-CoV-2', 'MERS-CoV'], 'Result': ['Positive', 'Negative'], 'Notes': ['Pass', 'Pass']})

def generate_lot_to_lot_data():
    df1 = pd.DataFrame({'Reagent Lot ID': 'Reference_Lot', 'Test Line Intensity': np.random.normal(100, 10, 30)})
    df2 = pd.DataFrame({'Reagent Lot ID': 'Candidate_Lot', 'Test Line Intensity': np.random.normal(102, 10, 30)})
    return pd.concat([df1, df2])

def calculate_equivalence(df, group_col, value_col, low, high):
    groups = df[group_col].unique(); data1 = df[df[group_col] == groups[0]][value_col]; data2 = df[df[group_col] == groups[1]][value_col]
    mean_diff = data2.mean() - data1.mean(); p_tost = 0.01 if low < mean_diff < high else 0.99
    return p_tost, 0.5, mean_diff

def generate_defect_trend_data():
    dates = pd.date_range(start="2023-01-01", periods=12, freq='W'); found = np.cumsum(np.random.randint(1, 5, 12))
    open_defects = found - np.cumsum(np.random.randint(0, 4, 12))
    return pd.DataFrame({'Date': dates, 'Total Defects Found': found, 'Open Defects': open_defects.clip(lower=0)})

def generate_submission_package_data(project_name, pathway):
    return pd.DataFrame({'Deliverable': ['V&V Master Plan', 'Risk Management File', 'RTM', 'Sensitivity Report', 'Specificity Report', 'V&V Summary Report'], 'Status': ['Approved', 'Approved', 'Approved', 'In Review', 'Execution', 'Drafting'], 'Progress': [100, 100, 100, 85, 40, 15], 'Regulatory_Impact': ['High', 'High', 'High', 'Medium', 'High', 'High'], 'Statistical_Robustness': [10, 9, 9, 7, 6, 9]})

def generate_capa_data():
    return pd.DataFrame({'ID': ['CAPA-23-011', 'INV-24-005', 'CAPA-24-002'], 'Source': ['Post-Launch Complaint', 'V&V Study Anomaly', 'Internal Audit'], 'Due Date': [date.today() - timedelta(days=10), date.today() + timedelta(days=18), date.today() + timedelta(days=40)]})

def generate_validation_program_data():
    data = {'System_ID': ['EQP-00123', 'EQP-00456', 'SW-0010', 'EQP-00789', 'UTL-0001'], 'System_Name': ['Automated Filler Line 3', 'Blister Pack Sealer 7', 'LIMS v3.2', 'Lyophilizer #2', 'WFI Loop'], 'Validation_Status': ['Validated', 'Validation in Progress', 'Validated', 'Revalidation Due', 'Validated'], 'Last_Validation_Date': [date(2022, 5, 20), date(2024, 6, 15), date(2023, 1, 10), date(2021, 8, 30), date(2024, 2, 28)], 'Next_Revalidation_Date': [date(2025, 5, 20), date(2027, 6, 15), date(2026, 1, 10), date(2024, 8, 30), date(2025, 2, 28)]}
    df = pd.DataFrame(data); df['Next_Revalidation_Date'] = pd.to_datetime(df['Next_Revalidation_Date'])
    df['Days_Until_Due'] = (df['Next_Revalidation_Date'] - pd.to_datetime(date.today())).dt.days
    df.loc[df['Days_Until_Due'] < 0, 'Validation_Status'] = 'Revalidation Overdue'
    return df

def generate_instrument_schedule_data():
    today=date.today(); return pd.DataFrame({'Instrument': ['Savanna-01', 'Savanna-01', 'Sofia-03', 'Vitros-02'], 'Start': [today, today+timedelta(days=2), today, today], 'Finish': [today + timedelta(days=1), today+timedelta(days=4), today + timedelta(days=5), today + timedelta(days=2)], 'Status': ['V&V Execution', 'Calibration/PM', 'V&V Execution', 'OOS']})

def generate_training_data_for_heatmap():
    return pd.DataFrame({'Savanna': [2, 1, 0, 1], 'Sofia': [1, 2, 1, 0], 'JMP/R': [2, 1, 1, 0]}, index=['M. Rodriguez', 'J. Chen', 'S. Patel', 'K. Lee'])

def generate_reagent_lot_status_data():
    return pd.DataFrame({'Material': ['RVP12 Master Mix', 'HIV Control Kit'], 'Lot': ['MM-24A01', 'CTRL-22C09'], 'Expiry': [date.today() + timedelta(days=180), date.today() - timedelta(days=10)], 'Status': ['Qualified', 'Expired']})

def calculate_instrument_utilization(schedule_df): return schedule_df

def generate_idp_data():
    return pd.DataFrame({'Team Member': ['J. Chen', 'S. Patel', 'K. Lee'], 'Development Goal': ['Expert in JMP/R', 'Practitioner on Vitros', 'Protocol Authoring'], 'Mentor': ['M. Rodriguez', 'Manager', 'M. Rodriguez'], 'Start Date': [date.today(), date.today(), date.today()], 'Target Date': [date.today() + timedelta(days=150), date.today() + timedelta(days=90), date.today() + timedelta(days=45)]})

def generate_process_excellence_data():
    months = pd.date_range(start="2023-01-01", periods=18, freq='M')
    return pd.DataFrame({'Month': months, 'Protocol_Approval_Cycle_Time_Days': 45 - np.arange(18) * 0.5 + np.random.normal(0, 2, 18), 'Report_Rework_Rate_Percent': 8 + np.sin(np.linspace(0, 2*np.pi, 18)) * 2, 'Deviations_per_100_Test_Hours': 3 - np.arange(18) * 0.05 + np.random.normal(0, 0.3, 18)})

def generate_workload_forecast_data():
    ds = pd.date_range(start='2022-01-01', periods=24, freq='MS'); trend = np.linspace(500, 800, 24)
    seasonality = 100 * np.sin(np.linspace(0, 4 * np.pi, 24)); y = trend + seasonality + np.random.normal(0, 50, 24)
    return pd.DataFrame({'ds': ds, 'y': y.clip(lower=100)})

def generate_monthly_review_ppt(kpi_data, fig_timeline, fig_risk):
    pres = Presentation(); pres.slide_width = Inches(10); pres.slide_height = Inches(7.5)
    title_slide = pres.slides.add_slide(pres.slide_layouts[0]); title_slide.shapes.title.text = "Validation Monthly Review"
    img_bytes = pio.to_image(fig_timeline, format="png", width=800, height=450, scale=2); pres.slides.add_slide(pres.slide_layouts[6]).shapes.add_picture(io.BytesIO(img_bytes), Inches(1), Inches(1.2), width=Inches(8))
    ppt_stream = io.BytesIO(); pres.save(ppt_stream); ppt_stream.seek(0)
    return ppt_stream
